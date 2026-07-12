import re
import json
import csv
import pandas as pd
from typing import List, Dict, Any, Tuple
from pathlib import Path

from langchain.text_splitter import RecursiveCharacterTextSplitter, CharacterTextSplitter
import tiktoken

from app.core.config import settings


class DocumentParser:
    """多类型文档解析器"""
    
    @staticmethod
    def parse_pdf(file_path: str) -> str:
        """解析 PDF 文件"""
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            content = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    content.append(text)
            return "\n\n".join(content)
        except Exception as e:
            raise Exception(f"解析 PDF 文件失败: {str(e)}")
    
    @staticmethod
    def parse_docx(file_path: str) -> str:
        """解析 DOCX 文件"""
        try:
            from docx import Document
            doc = Document(file_path)
            content = []
            for paragraph in doc.paragraphs:
                if paragraph.text:
                    content.append(paragraph.text)
            return "\n\n".join(content)
        except Exception as e:
            raise Exception(f"解析 DOCX 文件失败: {str(e)}")
    
    @staticmethod
    def parse_pptx(file_path: str) -> str:
        """解析 PPTX 文件"""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            content = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content.append(shape.text)
            return "\n\n".join(content)
        except Exception as e:
            raise Exception(f"解析 PPTX 文件失败: {str(e)}")
    
    @staticmethod
    def parse_excel(file_path: str) -> str:
        """解析 Excel 文件"""
        try:
            df = pd.read_excel(file_path)
            return df.to_csv(sep="\t", na_rep="", index=False)
        except Exception as e:
            raise Exception(f"解析 Excel 文件失败: {str(e)}")
    
    @staticmethod
    def parse_csv(file_path: str) -> str:
        """解析 CSV 文件"""
        try:
            df = pd.read_csv(file_path)
            return df.to_csv(sep="\t", na_rep="", index=False)
        except Exception as e:
            raise Exception(f"解析 CSV 文件失败: {str(e)}")
    
    @staticmethod
    def parse_text(file_path: str) -> str:
        """解析纯文本文件"""
        try:
            with open(file_path, encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            # 尝试 GBK 编码
            with open(file_path, encoding="gbk", errors="ignore") as f:
                return f.read()
        except Exception as e:
            raise Exception(f"解析文本文件失败: {str(e)}")
    
    @staticmethod
    def parse_json(file_path: str) -> str:
        """解析 JSON 文件"""
        try:
            with open(file_path, encoding="utf-8") as f:
                data = json.load(f)
            return json.dumps(data, ensure_ascii=False, indent=2)
        except Exception as e:
            raise Exception(f"解析 JSON 文件失败: {str(e)}")
    
    @classmethod
    def parse_file(cls, file_path: str, file_type: str) -> str:
        """根据文件类型解析文件"""
        if file_type == "pdf":
            return cls.parse_pdf(file_path)
        elif file_type == "docx":
            return cls.parse_docx(file_path)
        elif file_type == "pptx":
            return cls.parse_pptx(file_path)
        elif file_type == "excel":
            ext = Path(file_path).suffix.lower()
            if ext == ".csv":
                return cls.parse_csv(file_path)
            else:
                return cls.parse_excel(file_path)
        elif file_type in ["text", "structured", "other"]:
            return cls.parse_text(file_path)
        else:
            return cls.parse_text(file_path)


class TextCleaner:
    """文本清洗器"""
    
    @staticmethod
    def clean_text(text: str) -> str:
        """清洗文本内容"""
        if not text:
            return ""
        
        # 去除多余空白字符
        text = re.sub(r"\s+", " ", text)
        
        # 去除常见的无效字符
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
        
        # 规范化标点
        text = re.sub(r"\s+", " ", text)
        
        # 去除首尾空白
        text = text.strip()
        
        return text
    
    @staticmethod
    def remove_redundant_content(text: str) -> str:
        """去除多余的内容（页眉页脚、重复内容等）"""
        lines = text.split("\n")
        cleaned_lines = []
        seen_lines = set()
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 去除明显的页眉页脚（页码、日期等）
            if re.match(r"^\d+$", line) and len(line) <= 3:
                continue
            if re.match(r"^第\s*\d+\s*页", line):
                continue
            if re.match(r"^Page\s*\d+$", line, re.I):
                continue
            
            # 去除重复行（保留第一次出现）
            line_hash = hash(line)
            if line_hash in seen_lines and len(line) < 50:
                continue
            seen_lines.add(line_hash)
            
            cleaned_lines.append(line)
        
        return "\n".join(cleaned_lines)


class DocumentChunker:
    """文档切片器"""
    
    def __init__(self):
        self.tokenizer = tiktoken.get_encoding("cl100k_base")
        self.chunk_size = settings.CHUNK_SIZE
        self.chunk_overlap = settings.CHUNK_OVERLAP
    
    def count_tokens(self, text: str) -> int:
        """计算文本的 token 数量"""
        return len(self.tokenizer.encode(text))
    
    def chunk_text(
        self, 
        text: str, 
        chunk_size: Optional[int] = None, 
        chunk_overlap: Optional[int] = None
    ) -> List[Dict[str, Any]]:
        """
        将文本切分成片段
        
        返回:
            [{"content": "...", "tokens": ..., "start": ..., "end": ...}, ...]
        """
        chunk_size = chunk_size or self.chunk_size
        chunk_overlap = chunk_overlap or self.chunk_overlap
        
        if not text:
            return []
        
        # 使用 RecursiveCharacterTextSplitter 进行切分
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=self.count_tokens,
            separators=[
                "\n\n\n",
                "\n\n",
                "\n",
                ". ",
                "。",
                "! ",
                "！",
                "? ",
                "？",
                " ",
                "",
            ],
        )
        
        chunks = text_splitter.split_text(text)
        
        # 构建带元数据的片段
        result = []
        position = 0
        
        for idx, chunk in enumerate(chunks):
            tokens = self.count_tokens(chunk)
            start_pos = text.find(chunk, position)
            
            if start_pos == -1:
                start_pos = position
            
            end_pos = start_pos + len(chunk)
            position = end_pos
            
            result.append({
                "content": chunk,
                "tokens": tokens,
                "chunk_index": idx,
                "start_position": start_pos,
                "end_position": end_pos,
            })
        
        return result
    
    def chunk_by_paragraph(self, text: str) -> List[Dict[str, Any]]:
        """按段落切分"""
        paragraphs = re.split(r"\n\s*\n", text)
        
        result = []
        for idx, para in enumerate(paragraphs):
            para = para.strip()
            if para:
                tokens = self.count_tokens(para)
                result.append({
                    "content": para,
                    "tokens": tokens,
                    "chunk_index": idx,
                })
        
        return result


class BM25Preprocessor:
    """BM25 预处理工具"""
    
    def __init__(self):
        try:
            import nltk
            nltk.download("stopwords", quiet=True)
            nltk.download("punkt", quiet=True)
            from nltk.corpus import stopwords
            from nltk.tokenize import word_tokenize
            self.stopwords = set(stopwords.words("english") + ["的", "了", "在", "是", "我", "有", "和", "就"])
            self.word_tokenize = word_tokenize
        except Exception:
            self.stopwords = set()
            self.word_tokenize = None
    
    def tokenize(self, text: str) -> List[str]:
        """对文本进行分词（中英文混合）"""
        # 简单的中文分词策略
        # 在实际项目中建议使用 jieba 或 spaCy
        text = text.lower()
        
        # 保留中文汉字和英文单词
        pattern = r'[\u4e00-\u9fff]+|[a-zA-Z]+|\d+'
        tokens = re.findall(pattern, text)
        
        # 过滤停用词和短词
        tokens = [
            t for t in tokens 
            if t not in self.stopwords 
            and len(t) > 1
        ]
        
        return tokens
    
    def get_terms(self, text: str) -> List[str]:
        """获取文档的 BM25 关键词列表"""
        return self.tokenize(text)


# 初始化服务实例
document_parser = DocumentParser()
text_cleaner = TextCleaner()
document_chunker = DocumentChunker()
bm25_preprocessor = BM25Preprocessor()
